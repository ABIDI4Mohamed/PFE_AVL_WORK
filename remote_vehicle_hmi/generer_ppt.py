from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_slide(prs, title_text, bullet_points, image_placeholder_text=""):
    slide_layout = prs.slide_layouts[1] # Bullet point layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title = slide.shapes.title
    title.text = title_text
    
    # Corps
    tf = slide.placeholders[1].text_frame
    tf.text = bullet_points[0] if bullet_points else ""
    for bp in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = bp
        p.level = 0

    # Ajout d'une zone pour image si texte fourni
    if image_placeholder_text:
        left = Inches(5.5)
        top = Inches(2)
        width = Inches(4)
        height = Inches(4)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf_box = txBox.text_frame
        p = tf_box.add_paragraph()
        p.text = f"[{image_placeholder_text}]"
        p.alignment = PP_ALIGN.CENTER

def create_presentation():
    prs = Presentation()
    # Format 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1: Titre
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "État d’avancement PFE : Prototype de roulage à distance"
    slide.placeholders[1].text = "Conception d'une cabine de roulage assistée - AVL Maroc\nMohamed Abidi - 2026"

    # Slide 2: Architecture
    add_slide(prs, "Architecture Globale du Système", 
              ["Pilotage via HMI Python (PySide6)", 
               "Communication temps réel UDP", 
               "Supervision et calcul sous Simulink",
               "Feedback d'état vers l'opérateur"],
              "Insérer capture ARCHITECTURE")

    # Slide 3: Stateflow
    add_slide(prs, "Sécurisation : Logique Stateflow", 
              ["Positionnement 'Safety-First' avant la dynamique", 
               "Modes : Normal, Safe, Comm_loss, Fault, Emergency",
               "Gestion des transitions de sécurité"],
              "Insérer capture STATEFLOW")

    # Slide 4: Communication
    add_slide(prs, "Robustesse du lien UDP", 
              ["Fréquence Heartbeat : 50 Hz (20 ms)", 
               "Analyse RTT, Jitter et Checksum",
               "Supervision active du paquet (time_since_last_packet)"],
              "Insérer capture LOGIQUE RÉSEAU")

    # Slide 5: HMI
    add_slide(prs, "Interface Homme-Machine (HMI)", 
              ["Envoi : 11 signaux float32", 
               "Retour : 6 signaux float32",
               "Visualisation : Viewer 3D intégré"],
              "Insérer capture HMI")

    # Slide 6: Synthèse
    add_slide(prs, "Résultats et État Actuel", 
              ["Chaîne technique validée de bout en bout", 
               "Prototype cohérent et fonctionnel",
               "Enchaînement conforme aux objectifs du PFE"],
              "Insérer PHOTO SETUP / VIEWER")

    # Slide 7: Prochaines étapes
    add_slide(prs, "Perspectives et Optimisations", 
              ["Stabilisation des seuils Stateflow", 
               "Finalisation des comportements de sécurité",
               "Amélioration de l'ergonomie visuelle"],
              "Icônes WORK IN PROGRESS")

    # Slide 8: Conclusion
    add_slide(prs, "Conclusion", ["Base technique posée", "Objectifs respectés", "Questions ?"])

    # Sauvegarde
    filename = "presentation_avancement_pfe_avl.pptx"
    prs.save(filename)
    print(f"Présentation générée : {filename}")

if __name__ == "__main__":
    create_presentation()